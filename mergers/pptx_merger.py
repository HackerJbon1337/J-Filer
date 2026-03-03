"""
PPTX Merger Module
Merges multiple PowerPoint presentations into a single PPTX using python-pptx.
"""

import os
from pptx import Presentation
from copy import deepcopy
from lxml import etree


def merge_pptx(input_paths, output_path):
    """
    Merge multiple PPTX files into a single PowerPoint presentation.
    Copies slides from each source into a master presentation.

    Args:
        input_paths (list[str]): List of absolute paths to input PPTX files.
        output_path (str): Absolute path for the merged output PPTX.

    Returns:
        str: Path to the merged output file.

    Raises:
        ValueError: If fewer than 2 files are provided.
        FileNotFoundError: If any input file does not exist.
    """
    if len(input_paths) < 2:
        raise ValueError("At least 2 PPTX files are required for merging.")

    for path in input_paths:
        if not os.path.isfile(path):
            raise FileNotFoundError(f"Input file not found: {path}")

    # Use the first presentation as the base
    master_prs = Presentation(input_paths[0])

    for pptx_path in input_paths[1:]:
        source_prs = Presentation(pptx_path)
        _copy_slides(source_prs, master_prs)

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    master_prs.save(output_path)

    return output_path


def _copy_slides(source_prs, target_prs):
    """
    Copy all slides from source_prs into target_prs.
    Uses low-level XML manipulation to duplicate slide content.
    """
    for slide in source_prs.slides:
        # Add a new blank slide layout (use the first layout as fallback)
        slide_layout = target_prs.slide_layouts[6] if len(target_prs.slide_layouts) > 6 else target_prs.slide_layouts[0]
        new_slide = target_prs.slides.add_slide(slide_layout)

        # Clear default placeholders on new slide
        for placeholder in list(new_slide.placeholders):
            sp = placeholder._element
            sp.getparent().remove(sp)

        # Copy all shapes from source slide
        for shape in slide.shapes:
            el = deepcopy(shape.element)
            new_slide.shapes._spTree.append(el)

        # Copy slide background if present
        if slide.background.fill.type is not None:
            new_slide.background._element.attrib.update(
                slide.background._element.attrib
            )

    return target_prs
