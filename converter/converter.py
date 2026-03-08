"""
Cross-Format Converter Module
Supports all 6 conversions across PDF, DOCX, and PPTX:

  DOCX → PDF   (COM / Word)
  PPTX → PDF   (COM / PowerPoint)
  PDF  → DOCX  (PDF pages → images → python-docx)
  PDF  → PPTX  (PDF pages → images → python-pptx)
  DOCX → PPTX  (DOCX → PDF → PPTX pipeline)
  PPTX → DOCX  (PPTX → PDF → DOCX pipeline)

Requires Microsoft Office for any conversion that passes through Office formats.
PIL/Pillow and pikepdf are used for PDF→image rasterisation (no extra deps needed).
"""

import os
import tempfile


# ─── Conversion Matrix ────────────────────────────────────────────────────────
# Maps (source_ext, target_ext) → callable
_CONVERSION_MATRIX = {}   # populated after function definitions


def get_supported_conversions():
    """Return a list of supported (source_format, target_format) tuples."""
    return [
        ('pdf',  'docx'),
        ('pdf',  'pptx'),
        ('docx', 'pdf'),
        ('docx', 'pptx'),
        ('pptx', 'pdf'),
        ('pptx', 'docx'),
    ]


def convert_file(input_path, target_format, output_dir=None):
    """
    Unified entry-point for all supported conversions.

    Args:
        input_path   (str): Absolute path to the source file.
        target_format (str): Desired output format — 'pdf', 'docx', or 'pptx'.
        output_dir   (str | None): Directory for the converted file.
                                   Defaults to a system temp directory.

    Returns:
        str: Absolute path to the converted output file.

    Raises:
        ValueError:  Unsupported source/target combination or same format.
        RuntimeError: Conversion failed (missing Office, corrupt file, etc.).
    """
    input_path = os.path.abspath(input_path)
    src_ext = os.path.splitext(input_path)[1].lower().lstrip('.')

    # Normalise legacy extensions
    if src_ext == 'doc':
        src_ext = 'docx'
    elif src_ext == 'ppt':
        src_ext = 'pptx'

    target_format = target_format.lower().strip('.')

    if src_ext == target_format:
        raise ValueError(f"Source and target formats are the same: {target_format}")

    key = (src_ext, target_format)
    if key not in _CONVERSION_MATRIX:
        raise ValueError(
            f"Unsupported conversion: {src_ext.upper()} → {target_format.upper()}"
        )

    if output_dir is None:
        output_dir = tempfile.mkdtemp(prefix='jfiler_convert_')

    os.makedirs(output_dir, exist_ok=True)

    base_name = os.path.splitext(os.path.basename(input_path))[0]
    output_path = os.path.abspath(
        os.path.join(output_dir, f"{base_name}.{target_format}")
    )

    try:
        _CONVERSION_MATRIX[key](input_path, output_path)
    except RuntimeError:
        raise
    except Exception as e:
        raise RuntimeError(
            f"Conversion failed ({src_ext.upper()} → {target_format.upper()}): {e}"
        )

    return output_path


# ─── Legacy helper (keep backwards compatibility with app.py merger) ──────────
def convert_to_pdf(input_path, output_dir=None):
    """Convert DOCX/PPTX to PDF (backwards-compatible wrapper)."""
    return convert_file(input_path, 'pdf', output_dir)


# ─── Office Availability ──────────────────────────────────────────────────────
def is_office_available():
    """Return True if Microsoft Office COM objects are reachable."""
    try:
        import comtypes.client
        import pythoncom
        pythoncom.CoInitialize()
        try:
            word = comtypes.client.CreateObject('Word.Application')
            word.Quit()
            return True
        finally:
            pythoncom.CoUninitialize()
    except Exception:
        return False


# ─── Office → PDF ─────────────────────────────────────────────────────────────
def _docx_to_pdf(input_path, output_path):
    """DOCX (or DOC) → PDF via Word COM automation."""
    import comtypes.client
    import pythoncom

    pythoncom.CoInitialize()
    try:
        word = comtypes.client.CreateObject('Word.Application')
        word.Visible = False
        doc = word.Documents.Open(input_path)
        doc.SaveAs(output_path, FileFormat=17)   # 17 = wdFormatPDF
        doc.Close()
        word.Quit()
    finally:
        pythoncom.CoUninitialize()


def _pptx_to_pdf(input_path, output_path):
    """PPTX (or PPT) → PDF via PowerPoint COM automation."""
    import comtypes.client
    import pythoncom

    pythoncom.CoInitialize()
    try:
        ppt = comtypes.client.CreateObject('PowerPoint.Application')
        ppt.Visible = 1
        pres = ppt.Presentations.Open(input_path, WithWindow=False)
        pres.SaveAs(output_path, FileFormat=32)  # 32 = ppSaveAsPDF
        pres.Close()
        ppt.Quit()
    finally:
        pythoncom.CoUninitialize()


# ─── PDF → Images (shared rasteriser) ────────────────────────────────────────
def _pdf_to_images(pdf_path, dpi=150):
    """
    Rasterise each page of a PDF to a PIL Image object using PyMuPDF (fitz).

    Returns:
        list[PIL.Image.Image]: One image per PDF page (RGB mode).
    """
    import fitz  # PyMuPDF
    from PIL import Image as PILImage
    import io

    images = []
    zoom = dpi / 72.0
    matrix = fitz.Matrix(zoom, zoom)

    doc = fitz.open(pdf_path)
    for page in doc:
        pix = page.get_pixmap(matrix=matrix)
        img_data = pix.tobytes("png")
        img = PILImage.open(io.BytesIO(img_data)).convert('RGB')
        images.append(img)
    doc.close()

    return images


# ─── PDF → DOCX ───────────────────────────────────────────────────────────────
def _pdf_to_docx(input_path, output_path):
    """
    Convert PDF to DOCX using Word COM automation.
    Word can natively open PDFs and save as DOCX, preserving text and layout.
    """
    import comtypes.client
    import pythoncom

    pythoncom.CoInitialize()
    try:
        word = comtypes.client.CreateObject('Word.Application')
        word.Visible = False
        doc = word.Documents.Open(os.path.abspath(input_path))
        doc.SaveAs(os.path.abspath(output_path), FileFormat=16)  # 16 = wdFormatDocumentDefault (.docx)
        doc.Close()
        word.Quit()
    finally:
        pythoncom.CoUninitialize()


# ─── PDF → PPTX ───────────────────────────────────────────────────────────────
def _pdf_to_pptx(input_path, output_path):
    """
    Convert PDF to PPTX by creating one slide per page, each with the
    rasterised page image filling the entire slide area.
    """
    from pptx import Presentation
    from pptx.util import Inches, Emu
    import io

    images = _pdf_to_images(input_path, dpi=150)

    prs = Presentation()
    # Use blank slide layout
    blank_layout = prs.slide_layouts[6]

    # Set slide size to 10 × 7.5 inches (standard widescreen-ish)
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    for img in images:
        slide = prs.slides.add_slide(blank_layout)

        buf = io.BytesIO()
        img.save(buf, format='PNG')
        buf.seek(0)

        slide.shapes.add_picture(
            buf,
            left=Emu(0), top=Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(output_path)


# ─── Cross-format chained pipelines ──────────────────────────────────────────
def _docx_to_pptx(input_path, output_path):
    """DOCX → PDF (COM) → PPTX (image pipeline)."""
    tmp_dir = tempfile.mkdtemp(prefix='jfiler_chain_')
    pdf_path = os.path.join(tmp_dir, 'intermediate.pdf')
    _docx_to_pdf(input_path, pdf_path)
    _pdf_to_pptx(pdf_path, output_path)


def _pptx_to_docx(input_path, output_path):
    """PPTX → PDF (COM) → DOCX (image pipeline)."""
    tmp_dir = tempfile.mkdtemp(prefix='jfiler_chain_')
    pdf_path = os.path.join(tmp_dir, 'intermediate.pdf')
    _pptx_to_pdf(input_path, pdf_path)
    _pdf_to_docx(pdf_path, output_path)


# ─── Register all conversions ──────────────────────────────────────────────────
_CONVERSION_MATRIX = {
    ('docx', 'pdf'):  _docx_to_pdf,
    ('pptx', 'pdf'):  _pptx_to_pdf,
    ('pdf',  'docx'): _pdf_to_docx,
    ('pdf',  'pptx'): _pdf_to_pptx,
    ('docx', 'pptx'): _docx_to_pptx,
    ('pptx', 'docx'): _pptx_to_docx,
}
