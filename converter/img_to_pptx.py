"""
Image → PPTX Converter
Converts one or more image files (JPG, PNG, WEBP, BMP, TIFF, GIF)
into a single PowerPoint presentation (PPTX).

Each image is added to a new slide. High-resolution images are scaled 
down to fit the slide, while smaller images are kept at their original 
size or scaled up slightly to fit nicely.
"""

import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

def images_to_pptx(image_paths: list[str], output_path: str) -> None:
    """
    Merge one or more image files into a single PPTX.

    Args:
        image_paths (list[str]): Ordered list of absolute paths to image files.
        output_path (str):       Destination path for the output PPTX.
    """
    if not image_paths:
        raise ValueError("No images provided.")

    prs = Presentation()
    
    # Standard widespread slide format 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 6 is the index for a blank slide layout in the default template
    blank_slide_layout = prs.slide_layouts[6]

    for path in image_paths:
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Calculate scaling to fit the slide while preserving aspect ratio
        try:
            with Image.open(path) as img:
                img_width, img_height = img.size
        except Exception as e:
            continue
            
        # Slide dimensions in EMU (English Metric Units)
        # 1 Inch = 914400 EMU
        slide_w_emu = prs.slide_width
        slide_h_emu = prs.slide_height
        
        # Target bounding box (leave a small margin)
        margin = Inches(0)
        target_w = slide_w_emu - (2 * margin)
        target_h = slide_h_emu - (2 * margin)
        
        # Calculate aspect ratios
        slide_ratio = target_w / target_h
        img_ratio = img_width / img_height
        
        if img_ratio > slide_ratio:
            # Image is wider than slide: constrain by width
            final_w = target_w
            final_h = int(target_w / img_ratio)
        else:
            # Image is taller than slide: constrain by height
            final_h = target_h
            final_w = int(target_h * img_ratio)
            
        # Center the image
        left = int((slide_w_emu - final_w) / 2)
        top = int((slide_h_emu - final_h) / 2)
        
        try:
            slide.shapes.add_picture(
                os.path.abspath(path), 
                left, top, 
                width=final_w, 
                height=final_h
            )
        except Exception as e:
            # if add_picture fails (e.g. format not fully supported by pptx directly), 
            # we could try re-saving it as a temp PNG first
            import tempfile
            try:
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                    tmp_path = tmp.name
                with Image.open(path) as i:
                    i.convert('RGB').save(tmp_path, format='PNG')
                slide.shapes.add_picture(tmp_path, left, top, width=final_w, height=final_h)
                os.remove(tmp_path)
            except Exception as e2:
                pass

    prs.save(output_path)
